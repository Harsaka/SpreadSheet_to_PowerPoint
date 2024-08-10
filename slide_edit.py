import gspread
from oauth2client.service_account import ServiceAccountCredentials
from pptx import Presentation

# 1. スプレッドシートからデータを取得

# 認証とアクセスの設定
scope = ["https://spreadsheets.google.com/feeds", "https://www.googleapis.com/auth/drive"]
creds = ServiceAccountCredentials.from_json_keyfile_name('APIキーファイル名.json', scope)
client = gspread.authorize(creds)

# スプレッドシートを開く
sheet = client.open_by_key('問題データのスプシID').sheet1

# データを取得
data = sheet.get_all_records()

# 2. PowerPointのスライドにデータを埋め込む

# テンプレートファイルをを開く
prs = Presentation('template.pptx')

for row in data:
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)

    # プレースホルダーの文字列を変更
    slide.placeholders[0].text = row['スプシで問題文が格納されている列の名前(1行目)']
    slide.placeholders[1].text = row['スプシで答えが格納されている列の名前(1行目)']

# スライドをファイルに保存
prs.save('発表用スライド.pptx')
